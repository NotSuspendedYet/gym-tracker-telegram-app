#!/usr/bin/env python3
"""
Скрипт для создания презентации GymTracker для защиты в школе бизнеса.
Требования: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Цветовая схема (тёмная тема как в приложении)
COLORS = {
    'bg_primary': RGBColor(15, 15, 20),      # #0f0f14
    'bg_secondary': RGBColor(26, 26, 36),    # #1a1a24
    'bg_card': RGBColor(34, 34, 46),         # #22222e
    'accent': RGBColor(124, 58, 237),        # #7c3aed (фиолетовый)
    'accent_pink': RGBColor(236, 72, 153),   # #ec4899
    'text_primary': RGBColor(255, 255, 255), # белый
    'text_secondary': RGBColor(160, 160, 176), # #a0a0b0
    'success': RGBColor(34, 197, 94),        # #22c55e
    'warning': RGBColor(245, 158, 11),       # #f59e0b
    'danger': RGBColor(239, 68, 68),         # #ef4444
}


def set_slide_background(slide, color):
    """Устанавливает цвет фона слайда"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, top=0.3, font_size=44):
    """Добавляет заголовок слайда"""
    left = Inches(0.5)
    top = Inches(top)
    width = Inches(9)
    height = Inches(1)
    
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_primary']
    p.alignment = PP_ALIGN.LEFT
    
    return shape


def add_subtitle(slide, text, top=1.0):
    """Добавляет подзаголовок"""
    left = Inches(0.5)
    top = Inches(top)
    width = Inches(9)
    height = Inches(0.5)
    
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['text_secondary']
    
    return shape


def add_bullet_text(slide, items, left=0.5, top=1.5, width=9, font_size=18):
    """Добавляет текст с буллетами"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS['text_primary']
        p.space_after = Pt(12)
    
    return shape


def add_table(slide, data, left=0.5, top=1.8, col_widths=None):
    """Добавляет таблицу на слайд"""
    rows = len(data)
    cols = len(data[0]) if data else 0
    
    if col_widths is None:
        col_widths = [Inches(9 / cols)] * cols
    
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), 
                                    sum(col_widths), Inches(0.4 * rows)).table
    
    # Устанавливаем ширину колонок
    for i, width in enumerate(col_widths):
        table.columns[i].width = width
    
    # Заполняем данные
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            
            # Форматирование
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.alignment = PP_ALIGN.LEFT
            
            # Заголовок таблицы
            if row_idx == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLORS['text_primary']
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['accent']
            else:
                paragraph.font.color.rgb = COLORS['text_primary']
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['bg_card']
            
            # Вертикальное выравнивание
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return table


def create_presentation():
    """Создаёт презентацию"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Пустой layout
    blank_layout = prs.slide_layouts[6]
    
    # ========== СЛАЙД 0: Титульный ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    # Логотип/эмодзи
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🏋️ GymTracker"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # Подзаголовок
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Telegram-бот для отслеживания тренировок"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['text_secondary']
    p.alignment = PP_ALIGN.CENTER
    
    # Автор
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Презентация для защиты проекта"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text_secondary']
    p.alignment = PP_ALIGN.CENTER
    
    # ========== СЛАЙД 1: Почему хочу свой бизнес ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Почему я хочу свой бизнес?")
    
    items = [
        "Финансовая свобода — возможность самому распоряжаться временем",
        "Больше времени для себя и семьи",
        "Больше зарабатывать и влиять на свой заработок",
        "В найме (Android-разработчик) низкий потолок ЗП",
        "Хочу достичь большего и жить в лучших условиях",
        "Быть хозяином своей жизни, а не работать на чужие мечты"
    ]
    add_bullet_text(slide, items, top=1.5, font_size=22)
    
    # ========== СЛАЙД 2: Мой БОМЖ ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Мой БОМЖ", font_size=40)
    add_subtitle(slide, "Боли • Опасения • Мечты • Желания", top=0.9)
    
    bomzh_data = [
        ["", "💰 Финансы", "⏰ Время и свобода", "🎯 Самореализация"],
        ["Боли", 
         "Больше работаешь — столько же получаешь. Повышение редкое и небольшое",
         "Необходимость ходить в офис. Постоянный контроль работодателя",
         "Делаю не то, что хочу. Работаю на чужие цели"],
        ["Опасения",
         "Бедная старость из-за того, что всю жизнь работал в найме",
         "Работать в найме до старости без возможности отдыхать",
         "Прожить не ту жизнь, о которой мечтал"],
        ["Мечты",
         "Свой дом. Крутая машина или несколько машин",
         "Регулярные путешествия с семьёй",
         "Быть хозяином своей жизни"],
        ["Желания",
         "Больше денег прямо сейчас",
         "Побольше свободного времени",
         "Самому решать, что делать"]
    ]
    
    add_table(slide, bomzh_data, top=1.5, col_widths=[Inches(1.3), Inches(2.6), Inches(2.6), Inches(2.6)])
    
    # ========== СЛАЙД 3: БОМЖ клиентов ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "БОМЖ моих клиентов", font_size=40)
    
    clients_data = [
        ["", "🎓 Студент", "💼 Офисный работник", "👶 Молодая мама"],
        ["Боли", 
         "Выглядит слабее друзей, считает себя дрыщом. Полгода поднимает одни и те же веса",
         "Болит спина, постоянное похмелье, нет сил. В зеркале видит пивное пузо",
         "Не влезает в старую одежду. Нет времени на зал, не с кем оставить ребёнка"],
        ["Опасения",
         "Не накачает мышцы, останется дрыщом навсегда",
         "Здоровье ухудшится. Не знает как правильно заниматься",
         "Не вернёт форму после родов"],
        ["Мечты",
         "Нравиться девушкам, красивый внешний вид, спортивная форма",
         "Хорошая физическая форма, здоровье, энергия",
         "Вернуть форму как у подруги, которая быстро восстановилась"],
        ["Желания",
         "Понимать когда повышать веса, видеть свой прогресс",
         "Программа тренировок, знать что делать в зале",
         "Короткие тренировки дома (15 мин), пока ребёнок спит"]
    ]
    
    add_table(slide, clients_data, top=1.3, col_widths=[Inches(1.2), Inches(2.7), Inches(2.7), Inches(2.7)])
    
    # ========== СЛАЙД 4: ЦКП ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Ценный Конечный Продукт (ЦКП)", font_size=38)
    
    # Четыре блока
    ckp_items = [
        ("💎 ЦЕННЫЙ для покупателя", [
            "Видит свой прогресс в цифрах и графиках",
            "Мотивация продолжать тренировки",
            "Удобство — всё в Telegram, не нужно отдельное приложение",
            "Структурированные тренировки вместо хаоса"
        ]),
        ("💰 ЦЕННЫЙ для продавца", [
            "Подписочная модель — рекуррентный доход",
            "Низкая стоимость привлечения (Telegram-реклама)",
            "Масштабируемость без доп. затрат"
        ]),
        ("📊 КОНЕЧНЫЙ (измеримый результат)", [
            "История всех тренировок с датами",
            "Графики прогресса по весам и повторениям",
            "Статистика: сколько тренировок, какие мышцы качал"
        ]),
        ("📱 ПРОДУКТ (осязаемый результат)", [
            "Telegram Mini App с красивым интерфейсом",
            "Записи тренировок, которые не потеряются",
            "Личный кабинет с достижениями"
        ])
    ]
    
    top = 1.3
    for title, items in ckp_items:
        # Заголовок блока
        box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        
        top += 0.4
        
        # Пункты
        for item in items:
            box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(8.5), Inches(0.35))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_primary']
            top += 0.3
        
        top += 0.15
    
    # ========== СЛАЙД 5: Путь клиента (MVB) ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Путь клиента (MVB)", font_size=40)
    add_subtitle(slide, "От первого контакта до результата", top=0.9)
    
    steps = [
        ("1️⃣", "Первый контакт", "Увидел рекламу в Telegram-канале (спорт, ЗОЖ)"),
        ("2️⃣", "Интерес", "Нажал на рекламу, перешёл в бота, прочитал описание"),
        ("3️⃣", "Знакомство", "Открыл Mini App, посмотрел интерфейс и функции"),
        ("4️⃣", "Проба", "Добавил тестовую тренировку, убедился что работает"),
        ("5️⃣", "Использование", "Пошёл в зал, записал реальную тренировку — удобно!"),
        ("6️⃣", "Покупка", "Понял ценность, купил подписку (trial → платная)"),
        ("🎯", "Результат", "Записывает тренировки, видит прогресс, мотивирован расти")
    ]
    
    top = 1.5
    for emoji, title, desc in steps:
        # Эмодзи
        box = slide.shapes.add_textbox(Inches(0.3), Inches(top), Inches(0.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(20)
        
        # Заголовок шага
        box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(2), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        
        # Описание
        box = slide.shapes.add_textbox(Inches(2.8), Inches(top), Inches(6.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_primary']
        
        top += 0.55

    # ========== СЛАЙД 6: MVP (минимальный функционал) ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Минимальный функционал (MVP)", font_size=38)
    add_subtitle(slide, "Что нужно для запуска", top=0.9)
    
    mvp_data = [
        ["Шаг клиента", "Что нужно с моей стороны"],
        ["Первый контакт", "Простая реклама в 1-2 Telegram-каналах"],
        ["Интерес", "Бот с описанием и кнопкой 'Открыть приложение'"],
        ["Знакомство", "Mini App: главный экран, список упражнений"],
        ["Проба", "Возможность добавить упражнение и записать подход (вес × повторения)"],
        ["Использование", "История тренировок (список), базовый график прогресса"],
        ["Покупка", "Простая форма оплаты подписки (trial 7 дней → 99₽/мес)"],
        ["Результат", "Данные сохраняются, пользователь видит свою статистику"]
    ]
    
    add_table(slide, mvp_data, top=1.4, col_widths=[Inches(2.5), Inches(6.5)])
    
    # ========== СЛАЙД 7: Максимальный функционал ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Максимальный функционал", font_size=38)
    add_subtitle(slide, "Куда расти", top=0.9)
    
    max_data = [
        ["Шаг клиента", "Максимальные возможности"],
        ["Первый контакт", "Таргетированная реклама, блогеры, несколько каналов, SEO"],
        ["Интерес", "Видео-превью, отзывы пользователей, FAQ в боте"],
        ["Знакомство", "Красивый UI, разные типы упражнений (силовые, кардио, плавание, статика)"],
        ["Проба", "Онбординг с подсказками, готовые программы тренировок"],
        ["Использование", "Детальная статистика, достижения, напоминания, экспорт данных"],
        ["Покупка", "Несколько тарифов, семейная подписка, годовая со скидкой"],
        ["Результат", "Социальные функции: делиться прогрессом, челленджи с друзьями, рейтинги"]
    ]
    
    add_table(slide, max_data, top=1.4, col_widths=[Inches(2.5), Inches(6.5)])
    
    # ========== СЛАЙД 8: Конкуренты ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Анализ конкурентов", font_size=40)
    
    competitors_data = [
        ["Конкурент", "Плюсы", "Минусы", "Цена"],
        ["gym-bot.ru", "Красивый, планы питания, анализ", "Нет интерфейса (только бот)", "500₽/мес"],
        ["@gym_ru_bot", "Крутые программы, бесплатный", "Нет интерфейса, нет статистики", "Бесплатно"],
        ["Gym Tracker (РуСтор)", "Гибкая статистика", "Некрасивый, надо заходить самому", "Бесплатно"],
        ["50 подтягиваний и т.п.", "Программы роста, достижения", "Узкая специализация, реклама", "Бесплатно"],
        ["Тетрадка", "Дёшево, тактильно", "Неудобно, можно потерять", "~50₽"],
        ["🏋️ GymTracker (я)", "Удобно в Telegram, красиво, статистика", "Пока MVP", "99₽/мес"]
    ]
    
    add_table(slide, competitors_data, top=1.3, col_widths=[Inches(2.2), Inches(2.8), Inches(2.8), Inches(1.3)])
    
    # ========== СЛАЙД 9: УТП ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Уникальное Торговое Предложение", font_size=36)
    
    # Главный месседж
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "«Ходишь в зал? А какой вес брал на тяге месяц назад? Не помнишь? Нет роста веса — нет роста мышц. Записывай тренировки в боте и смотри как растёшь!»"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    utp_items = [
        "✅ Удобно — прямо в Telegram, не нужно отдельное приложение",
        "✅ Красиво — современный тёмный интерфейс как в топовых приложениях",
        "✅ Функционально — разные типы тренировок, графики, история",
        "✅ Мотивирует — видишь свой прогресс, хочется расти дальше",
        "✅ Социально — делись успехами с друзьями, парные челленджи"
    ]
    
    add_bullet_text(slide, utp_items, top=2.7, font_size=20)
    
    # ========== СЛАЙД 10: Итоги ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    add_title_shape(slide, "Итоги и следующие шаги", font_size=40)
    
    summary_items = [
        "✅ Проблема понятна: люди не отслеживают прогресс → нет роста",
        "✅ Решение есть: удобный бот в Telegram с красивым интерфейсом",
        "✅ MVP готов: можно записывать тренировки и видеть статистику",
        "✅ Монетизация: подписочная модель (trial + 99₽/мес)",
        "",
        "📍 Следующие шаги:",
        "   1. Доработать функционал подписки",
        "   2. Запустить рекламу в 2-3 спортивных каналах",
        "   3. Собрать первых 100 пользователей",
        "   4. Получить обратную связь и улучшить продукт"
    ]
    
    add_bullet_text(slide, summary_items, top=1.3, font_size=20)
    
    # ========== СЛАЙД 11: Спасибо ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['bg_primary'])
    
    # Логотип
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "🏋️ GymTracker"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # Спасибо
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Спасибо за внимание!"
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['text_primary']
    p.alignment = PP_ALIGN.CENTER
    
    # Контакт
    box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Вопросы?"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['text_secondary']
    p.alignment = PP_ALIGN.CENTER
    
    # Сохраняем
    output_path = "GymTracker_Presentation.pptx"
    prs.save(output_path)
    print(f"[OK] Presentation saved: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
